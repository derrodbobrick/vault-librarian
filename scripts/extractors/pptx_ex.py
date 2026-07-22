"""PowerPoint (.pptx) extraction: per-slide text, notes, tables, images, renders.

Slides are layout- and chart-heavy, so the full-page render of each slide is
the primary signal; the text/notes give exact wording and speaker intent.
"""
from __future__ import annotations

from .common import Bundle
from . import render


def extract(src: str, out_dir: str, max_pages: int = 0) -> dict:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    b = Bundle(src, out_dir, "pptx")
    prs = Presentation(src)

    cp = prs.core_properties
    b.meta = {
        "title": cp.title or "",
        "author": cp.author or "",
        "created": str(cp.created or ""),
        "slideCount": len(prs.slides),
    }
    b.add(f"# PowerPoint: {cp.title or src.split('/')[-1]}")
    b.add(f"Slides: {len(prs.slides)}")
    if cp.author:
        b.add(f"Author: {cp.author}")

    charts = 0
    for idx, slide in enumerate(prs.slides, 1):
        b.heading(f"Slide {idx}", 2)
        for shape in slide.shapes:
            try:
                if shape.has_chart:
                    charts += 1
                    ch = shape.chart
                    title = ""
                    try:
                        if ch.has_title:
                            title = ch.chart_title.text_frame.text
                    except Exception:
                        pass
                    b.add(f"[chart: {ch.chart_type}{' — ' + title if title else ''} "
                          "→ see the slide render]")
                    continue
            except Exception:
                pass
            try:
                if shape.has_table:
                    tbl = shape.table
                    for ri, row in enumerate(tbl.rows):
                        cells = [c.text.replace("\n", " ").replace("|", "/") for c in row.cells]
                        b.add("| " + " | ".join(cells) + " |")
                        if ri == 0:
                            b.add("| " + " | ".join(["---"] * len(cells)) + " |")
                    continue
            except Exception:
                pass
            if shape.has_text_frame and shape.text_frame.text.strip():
                b.add(shape.text_frame.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    b.add_media(img.blob, img.ext or "png", desc=f"image on slide {idx}")
                except Exception:
                    pass

        # Speaker notes carry the presenter's intent.
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    b.add(f"*Speaker notes:* {notes}")
        except Exception:
            pass

    b.meta["charts"] = charts
    b.meta["embeddedImages"] = len(b.media)
    render.render_office_pages(src, b, max_pages)
    return b.result()
